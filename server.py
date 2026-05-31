from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import requests as req
import pickle
import os
from dotenv import load_dotenv
load_dotenv()
import json
import uuid
import io
import time
import base64
import threading
import logging
logging.basicConfig(level=logging.INFO)
from datetime import datetime, timedelta
from PIL import Image
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
import google.auth.transport.requests
import google.oauth2.service_account

app = Flask(__name__)

YOOKASSA_SHOP_ID = os.getenv("YOOKASSA_SHOP_ID")
YOOKASSA_SECRET_KEY = os.getenv("YOOKASSA_SECRET_KEY")
CORS(app)

@app.after_request
def after_request(response):
    response.headers.add('Access-Control-Allow-Origin', '*')
    response.headers.add('Access-Control-Allow-Headers', 'Content-Type,X-API-Token,X-Session-Token')
    response.headers.add('Access-Control-Allow-Methods', 'GET,POST,OPTIONS,DELETE')
    response.headers.add('Access-Control-Allow-Headers', 'Content-Type,X-API-Token,X-Session-Token,Accept,Origin')
    return response

BOT_TOKEN = os.getenv("BOT_TOKEN")
COMFYUI_URL = "http://127.0.0.1:8188"
API_TOKEN = os.getenv("API_TOKEN")
BASE_DIR = "E:/ComfyUI_windows_portable"

AUTH_CODES_FILE = os.path.join(BASE_DIR, "auth_codes.pkl")
TOKENS_FILE = os.path.join(BASE_DIR, "push_tokens.pkl")
WEB_SESSIONS_FILE = os.path.join(BASE_DIR, "web_sessions.pkl")
WEB_V4FLASH_HISTORY_FILE = os.path.join(BASE_DIR, "web_v4flash_history.pkl")
WEB_V4PRO_HISTORY_FILE = os.path.join(BASE_DIR, "web_v4pro_history.pkl")
V4FLASH_BALANCE_FILE = os.path.join(BASE_DIR, "v4flash_balance.pkl")
V4PRO_BALANCE_FILE = os.path.join(BASE_DIR, "v4pro_balance.pkl")
VIP_FILE = os.path.join(BASE_DIR, "vip_users.pkl")
VIP_EXPIRY_FILE = os.path.join(BASE_DIR, "vip_expiry.pkl")
FREE_GENERATIONS_FILE = os.path.join(BASE_DIR, "free_generations.pkl")
FREE_CHAT_HISTORY_FILE = os.path.join(BASE_DIR, "bot_site_free_history.pkl")
UNVIEWED_FILE = os.path.join(BASE_DIR, "unviewed_files.pkl")

VIP_USERS = set()
VIP_EXPIRY = {}
free_generations = {}
unviewed_files = {}

if os.path.exists(VIP_FILE):
    try:
        VIP_USERS = set(pickle.load(open(VIP_FILE, "rb")))
    except:
        pass
if os.path.exists(VIP_EXPIRY_FILE):
    try:
        VIP_EXPIRY = pickle.load(open(VIP_EXPIRY_FILE, "rb"))
    except:
        pass
if os.path.exists(FREE_GENERATIONS_FILE):
    try:
        free_generations = pickle.load(open(FREE_GENERATIONS_FILE, "rb"))
    except:
        pass
if os.path.exists(UNVIEWED_FILE):
    try:
        unviewed_files = pickle.load(open(UNVIEWED_FILE, "rb"))
    except:
        pass

auth_codes = {}
web_sessions = {}
web_sessions_time = {}
push_tokens = {}
web_v4flash_history = {}
web_v4pro_history = {}
user_v4flash_balance = {}
user_v4pro_balance = {}

TEMP_FOLDER = os.path.join(BASE_DIR, "temp_previews")
DATASET_FOLDER = os.path.join(BASE_DIR, "ComfyUI/input/my_first_lora_dataset")
INPUT_FOLDER = os.path.join(BASE_DIR, "ComfyUI/input")
WORKFLOWS_DIR = os.path.join(BASE_DIR, "ComfyUI/user/default/workflows/workflows")
OUTPUT_FOLDER = os.path.join(BASE_DIR, "ComfyUI/output")

os.makedirs(TEMP_FOLDER, exist_ok=True)
os.makedirs(DATASET_FOLDER, exist_ok=True)

SERVICE_ACCOUNT_FILE = os.path.join(BASE_DIR, "firebase_key.json")
firebase_credentials = google.oauth2.service_account.Credentials.from_service_account_file(
    SERVICE_ACCOUNT_FILE,
    scopes=["https://www.googleapis.com/auth/firebase.messaging"]
)

app_jobs = {}
app_jobs_lock = threading.Lock()
rate_limit = {}
rate_limit_lock = threading.Lock()

def save_vip():
    temp_vip = VIP_FILE + ".tmp"
    temp_expiry = VIP_EXPIRY_FILE + ".tmp"
    try:
        with open(temp_vip, "wb") as f:
            pickle.dump(list(VIP_USERS), f)
        os.replace(temp_vip, VIP_FILE)
        with open(temp_expiry, "wb") as f:
            pickle.dump(VIP_EXPIRY, f)
        os.replace(temp_expiry, VIP_EXPIRY_FILE)
    except Exception as e:
        print(f"[SAVE VIP] Error: {e}")

def safe_pickle_dump(data, filepath):
    temp_path = filepath + ".tmp"
    try:
        with open(temp_path, "wb") as f:
            pickle.dump(data, f)
        os.replace(temp_path, filepath)
    except Exception as e:
        print(f"[SAVE] Error saving {filepath}: {e}")

def check_rate_limit(ip, endpoint, max_requests, period_seconds):
    now = time.time()
    with rate_limit_lock:
        if ip not in rate_limit:
            rate_limit[ip] = {}
        if endpoint not in rate_limit[ip]:
            rate_limit[ip][endpoint] = []
        rate_limit[ip][endpoint] = [t for t in rate_limit[ip][endpoint] if now - t < period_seconds]
        if len(rate_limit[ip][endpoint]) >= max_requests:
            return False
        rate_limit[ip][endpoint].append(now)
        return True

def check_api_token(request):
    token = request.headers.get("X-API-Token") or request.form.get("token") or request.args.get("token")
    return token == API_TOKEN

def check_web_session(request):
    token = request.headers.get("X-Session-Token") or request.args.get("session_token")
    if not token:
        return False
    if token in web_sessions:
        if token in web_sessions_time:
            if datetime.now() - web_sessions_time[token] > timedelta(hours=24):
                del web_sessions[token]
                del web_sessions_time[token]
                safe_pickle_dump((web_sessions, web_sessions_time), WEB_SESSIONS_FILE)
                return False
        return True
    return False

def save_auth_codes():
    safe_pickle_dump(auth_codes, AUTH_CODES_FILE)

def load_all():
    global push_tokens, web_sessions, web_v4flash_history, web_v4pro_history, user_v4flash_balance, user_v4pro_balance
    if os.path.exists(TOKENS_FILE):
        with open(TOKENS_FILE, "rb") as f:
            push_tokens = pickle.load(f)
    if os.path.exists(WEB_SESSIONS_FILE):
        with open(WEB_SESSIONS_FILE, "rb") as f:
            data = pickle.load(f)
            if isinstance(data, tuple):
                web_sessions, web_sessions_time = data
            else:
                web_sessions = data
    if os.path.exists(WEB_V4FLASH_HISTORY_FILE):
        with open(WEB_V4FLASH_HISTORY_FILE, "rb") as f:
            web_v4flash_history = pickle.load(f)
    if os.path.exists(WEB_V4PRO_HISTORY_FILE):
        with open(WEB_V4PRO_HISTORY_FILE, "rb") as f:
            web_v4pro_history = pickle.load(f)
    if os.path.exists(V4FLASH_BALANCE_FILE):
        with open(V4FLASH_BALANCE_FILE, "rb") as f:
            user_v4flash_balance = pickle.load(f)
    if os.path.exists(V4PRO_BALANCE_FILE):
        with open(V4PRO_BALANCE_FILE, "rb") as f:
            user_v4pro_balance = pickle.load(f)

def save_tokens():
    safe_pickle_dump(push_tokens, TOKENS_FILE)

def load_free_history():
    global free_chat_history
    if os.path.exists(FREE_CHAT_HISTORY_FILE):
        with open(FREE_CHAT_HISTORY_FILE, "rb") as f:
            free_chat_history = pickle.load(f)

def send_push(user_id, title, body):
    if user_id not in push_tokens:
        print(f"[PUSH] No token for user {user_id}")
        return
    token = push_tokens[user_id]
    google.auth.transport.requests.Request()
    firebase_credentials.refresh(google.auth.transport.requests.Request())
    access_token = firebase_credentials.token
    url = f"https://fcm.googleapis.com/v1/projects/photobot-4309a/messages:send"
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    payload = {
        "message": {
            "token": token,
            "notification": {
                "title": title,
                "body": body
            }
        }
    }
    resp = req.post(url, headers=headers, json=payload)
    print(f"[PUSH] Sent to {user_id}: {resp.status_code} {resp.text}")

load_all()

def save_image_to_input(img_data, image_name):
    path = os.path.join(INPUT_FOLDER, image_name)
    with open(path, "wb") as f:
        f.write(img_data)
    return image_name

def load_workflow(mode, image_name, caption=""):
    path = os.path.join(WORKFLOWS_DIR, f"{mode}.json")
    if not os.path.exists(path):
        return None
    with open(path, "r", encoding="utf-8") as f:
        wf = json.load(f)
    wf_str = json.dumps(wf)
    if image_name:
        wf_str = wf_str.replace("{{image_name}}", image_name)
    wf_str = wf_str.replace("{{caption}}", caption)
    wf_str = wf_str.replace("{{clothes}}", caption)
    wf_str = wf_str.replace("{{dress}}", caption)
    wf_str = wf_str.replace("{{age}}", caption)
    wf_str = wf_str.replace("{{noise_seed}}", str(uuid.uuid4().int >> 64))
    return json.loads(wf_str)

# === ЭНДПОИНТЫ ===

@app.route("/check_subscription", methods=["POST"])
def check_subscription():
    data = request.json
    user_id = data.get("user_id")
    if not user_id:
        return jsonify({"subscribed": False, "error": "No user_id"})
    url = f"https://api.telegram.org/bot{BOT_TOKEN}/getChatMember"
    params = {"chat_id": user_id, "user_id": user_id}
    try:
        resp = req.get(url, params=params, timeout=5).json()
        if resp.get("ok"):
            status = resp["result"]["status"]
            return jsonify({"subscribed": status in ("member", "administrator", "creator")})
    except:
        pass
    return jsonify({"subscribed": False, "error": "API error"})

@app.route("/check_app_subscription", methods=["GET"])
def check_app_subscription():
    return jsonify({"subscribed": True})

@app.route("/verify_code", methods=["POST"])
def verify_code():
    ip = request.headers.get("X-Forwarded-For", request.remote_addr)
    if not check_rate_limit(ip, "verify_code", 5, 60):
        return jsonify({"success": False, "error": "Слишком много попыток. Подождите минуту."})
    global auth_codes, web_sessions
    if os.path.exists(AUTH_CODES_FILE):
        with open(AUTH_CODES_FILE, "rb") as f:
            auth_codes = pickle.load(f)
    data = request.json
    code = data.get("code")
    if not code:
        return jsonify({"success": False, "error": "No code"})
    user_id = None
    for uid, c in auth_codes.items():
        if c == code:
            user_id = uid
            break
    if not user_id:
        return jsonify({"success": False, "error": "Invalid code"})
    username = None
    try:
        chat_resp = req.get(f"https://api.telegram.org/bot{BOT_TOKEN}/getChat?chat_id={user_id}", timeout=5).json()
        if chat_resp.get("ok"):
            username = chat_resp["result"].get("username")
    except:
        pass
    url = f"https://api.telegram.org/bot{BOT_TOKEN}/getChatMember"
    params = {"chat_id": user_id, "user_id": user_id}
    try:
        resp = req.get(url, params=params, timeout=5).json()
        if resp.get("ok"):
            status = resp["result"]["status"]
            if status in ("member", "administrator", "creator"):
                del auth_codes[user_id]
                save_auth_codes()
                session_token = uuid.uuid4().hex
                web_sessions[session_token] = user_id
                web_sessions_time[session_token] = datetime.now()
                safe_pickle_dump((web_sessions, web_sessions_time), WEB_SESSIONS_FILE)
                return jsonify({"success": True, "user_id": user_id, "username": username, "session_token": session_token, "api_token": API_TOKEN, "subscribed": True})
            return jsonify({"success": False, "error": "Not subscribed", "user_id": user_id, "username": username})
    except:
        pass
    return jsonify({"success": False, "error": "API error"})

@app.route("/generate", methods=["POST"])
def generate():
    if not check_api_token(request):
        return jsonify({"error": "Unauthorized"}), 403
    mode = request.form.get('mode', 'documents')
    caption = request.form.get('caption', '')
    image_name = None
    if mode != "wallpaper":
        if 'image' not in request.files:
            return jsonify({"error": "No image"}), 400
        file = request.files['image']
        img_data = file.read()
        image_name = f"api_{uuid.uuid4().hex[:8]}.jpg"
        save_image_to_input(img_data, image_name)
    caption_val = caption.strip().replace('\n', ' ').replace('\r', ' ').replace('"', "'") if caption.strip() else ""
    if mode == "documents":
        caption_val = caption_val if caption_val else "черная рубашка"
    elif mode == "birthday":
        parts = caption_val.split() if caption_val else []
        caption_val = " ".join(parts) if parts else "23 красное"
    elif mode == "flowers":
        caption_val = caption_val if caption_val else "белое"
    wf = load_workflow(mode, image_name, caption_val)
    if wf is None:
        return jsonify({"error": "Workflow not found"}), 400
    prompt_resp = req.post(f"{COMFYUI_URL}/prompt", json={"prompt": wf}).json()
    prompt_id = prompt_resp.get("prompt_id")
    if not prompt_id:
        return jsonify({"error": "Prompt failed"}), 500
    for i in range(180):
        time.sleep(1)
        history = req.get(f"{COMFYUI_URL}/history").json()
        if prompt_id in history:
            outputs = history[prompt_id]["outputs"]
            for nid in outputs:
                for img in outputs[nid].get("images", []):
                    if img["filename"].startswith("Flux2-Klein"):
                        img_resp = req.get(f"{COMFYUI_URL}/view?filename={img['filename']}&subfolder={img.get('subfolder','')}&type=output")
                        bio = io.BytesIO(img_resp.content)
                        bio.seek(0)
                        return send_file(bio, mimetype="image/jpeg")
    return jsonify({"error": "Timeout"}), 500

@app.route("/enqueue", methods=["POST"])
def enqueue():
    source = request.form.get('source', 'app')
    if source == 'web':
        if not check_web_session(request):
            return jsonify({"error": "Unauthorized"}), 403
    elif not check_api_token(request):
        return jsonify({"error": "Unauthorized"}), 403
    mode = request.form.get('mode', 'documents')
    user_id = request.form.get('user_id', 'unknown')
    caption = request.form.get('caption', '')
    ref_name = request.form.get('ref_name', '')
    img_data = None
    if mode != "wallpaper":
        if 'image' not in request.files and not ref_name:
            return jsonify({"error": "No image"}), 400
        if 'image' in request.files:
            file = request.files['image']
            img_data = file.read()
    job_id = f"app_{user_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:6]}"
    print(f"[ENQUEUE] {job_id} for user {user_id} mode {mode} ref {ref_name}")
    with app_jobs_lock:
        app_jobs[job_id] = {
            "user_id": user_id,
            "status": "pending",
            "image_data": img_data,
            "caption": caption,
            "mode": mode,
            "size": request.form.get('size', 'pc'),
            "ref_name": ref_name,
            "source": source,
        }
    print(f"[ENQUEUE] Queue size: {len(app_jobs)}")
    return jsonify({"success": True, "job_id": job_id})

@app.route("/upload_reference", methods=["POST"])
def upload_reference():
    if not check_api_token(request):
        return jsonify({"error": "Unauthorized"}), 403
    if 'image' not in request.files:
        return jsonify({"error": "No image"}), 400
    user_id = request.form.get('user_id', 'unknown')
    unique_id = request.form.get('unique_id', '')
    if not unique_id:
        unique_id = f"user_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:6]}"
    file = request.files['image']
    user_folder = os.path.join(INPUT_FOLDER, f"user_{user_id}", unique_id)
    os.makedirs(user_folder, exist_ok=True)
    file.save(os.path.join(user_folder, "user_photo.jpg"))
    ref_name = f"user_{user_id}/{unique_id}/user_photo.jpg"
    return jsonify({"success": True, "ref_name": ref_name, "unique_id": unique_id})

@app.route("/upload_etalon", methods=["POST"])
def upload_etalon():
    if not check_api_token(request):
        return jsonify({"error": "Unauthorized"}), 403
    user_id = request.form.get('user_id', 'unknown')
    if 'image' not in request.files:
        return jsonify({"error": "No image"}), 400
    file = request.files['image']
    user_folder = os.path.join(INPUT_FOLDER, f"user_{user_id}")
    os.makedirs(user_folder, exist_ok=True)
    file.save(os.path.join(user_folder, "etalon.jpg"))
    return jsonify({"success": True})

@app.route("/check_etalon", methods=["GET"])
def check_etalon():
    user_id = request.args.get("user_id", "")
    etalon_path = os.path.join(INPUT_FOLDER, f"user_{user_id}", "etalon.jpg")
    return jsonify({"has_etalon": os.path.exists(etalon_path)})

@app.route("/get_etalon", methods=["GET"])
def get_etalon():
    user_id = request.args.get("user_id", "")
    etalon_path = os.path.join(INPUT_FOLDER, f"user_{user_id}", "etalon.jpg")
    if os.path.exists(etalon_path):
        return send_file(etalon_path, mimetype="image/jpeg")
    return jsonify({"error": "Not found"}), 404

@app.route("/dataset", methods=["POST"])
def dataset():
    if not check_api_token(request):
        return jsonify({"error": "Unauthorized"}), 403
    if 'image' not in request.files:
        return jsonify({"error": "No image"}), 400
    file = request.files['image']
    prompt = request.form.get('prompt', '')
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    uid = uuid.uuid4().hex[:6]
    fname = f"app_{ts}_{uid}"
    file.save(os.path.join(DATASET_FOLDER, fname + "_output.jpg"))
    with open(os.path.join(DATASET_FOLDER, fname + "_output.txt"), "w", encoding="utf-8") as f:
        f.write(prompt)
    return jsonify({"success": True})

@app.route("/pending", methods=["GET"])
def pending():
    source = request.args.get("source", "app")
    if source == 'web':
        if not check_web_session(request):
            return jsonify({"error": "Unauthorized"}), 403
    elif not check_api_token(request):
        return jsonify({"error": "Unauthorized"}), 403
    user_id = request.args.get("user_id")
    after = request.args.get("after", "")
    print(f"[PENDING] user={user_id} after={after}")
    results = []
    with app_jobs_lock:
        for job_id, job in list(app_jobs.items()):
            if job.get("user_id") == user_id and job.get("status") == "ready" and job.get("source") == source:
                if job_id > after:
                    if "file" in job and os.path.exists(job["file"]):
                        with open(job["file"], "rb") as f:
                            b64 = base64.b64encode(f.read()).decode()
                        results.append({"job_id": job_id, "data": b64, "ref_name": job.get("ref_name", "")})
                        print(f"[PENDING] Returning {job_id}")
                        # Добавляем файл в непросмотренные
                        user_id_str = str(user_id)
                        if user_id_str not in unviewed_files:
                            unviewed_files[user_id_str] = []
                        unviewed_files[user_id_str].append(job_id)
                        safe_pickle_dump(unviewed_files, UNVIEWED_FILE)
                        del app_jobs[job_id]
    print(f"[PENDING] Results: {len(results)}")
    return jsonify({"results": results})

@app.route("/unviewed_count", methods=["GET"])
def unviewed_count():
    source = request.args.get("source", "web")
    if source == 'web':
        if not check_web_session(request):
            return jsonify({"error": "Unauthorized"}), 403
    elif not check_api_token(request):
        return jsonify({"error": "Unauthorized"}), 403
    user_id = request.args.get("user_id", "")
    user_id_str = str(user_id)
    count = len(unviewed_files.get(user_id_str, []))
    return jsonify({"count": count})

@app.route("/mark_viewed", methods=["POST"])
def mark_viewed():
    source = request.args.get("source", "web")
    if source == 'web':
        if not check_web_session(request):
            return jsonify({"error": "Unauthorized"}), 403
    elif not check_api_token(request):
        return jsonify({"error": "Unauthorized"}), 403
    user_id = request.args.get("user_id", "")
    user_id_str = str(user_id)
    unviewed_files[user_id_str] = []
    safe_pickle_dump(unviewed_files, UNVIEWED_FILE)
    return jsonify({"success": True})

@app.route("/history_file", methods=["GET"])
def history_file():
    source = request.args.get("source", "app")
    if source == 'web':
        if not check_web_session(request):
            return jsonify({"error": "Unauthorized"}), 403
    elif not check_api_token(request):
        return jsonify({"error": "Unauthorized"}), 403
    filename = request.args.get("file", "")
    path = os.path.join(DATASET_FOLDER, filename)
    if os.path.exists(path):
        return send_file(path, mimetype="image/jpeg")
    return jsonify({"error": "Not found"}), 404

@app.route("/history", methods=["GET"])
def history():
    source = request.args.get("source", "app")
    if source == 'web':
        if not check_web_session(request):
            return jsonify({"error": "Unauthorized"}), 403
    elif not check_api_token(request):
        return jsonify({"error": "Unauthorized"}), 403
    user_id = request.args.get("user_id", "")
    files = []
    for f in os.listdir(DATASET_FOLDER):
        if f.startswith(f"{source}_{user_id}_") and f.endswith(".jpg"):
            files.append(f)
    files.sort(reverse=True)
    return jsonify({"files": files[:50]})

@app.route("/active", methods=["GET"])
def active():
    user_id = request.args.get("user_id", "")
    source = request.args.get("source", "")
    with app_jobs_lock:
        for job_id, job in app_jobs.items():
            if job.get("user_id") == user_id and job.get("status") in ("pending", "ready") and job.get("source") == source:
                return jsonify({"active": True, "job_id": job_id})
    return jsonify({"active": False, "job_id": ""})

@app.route("/register_token", methods=["POST"])
def register_token():
    if not check_api_token(request):
        return jsonify({"error": "Unauthorized"}), 403
    global push_tokens
    data = request.json
    user_id = data.get("user_id")
    token = data.get("token")
    push_tokens[user_id] = token
    save_tokens()
    print(f"[TOKEN] User {user_id}: {token}")
    return jsonify({"success": True})

@app.route("/get_full_response", methods=["GET"])
def get_full_response():
    if not check_web_session(request):
        return jsonify({"error": "Unauthorized"}), 403
    user_id_str = request.args.get("user_id", "")
    model = request.args.get("model", "flash")
    chat_id = request.args.get("chat_id", "default")
    user_id_int = int(user_id_str) if user_id_str.isdigit() else user_id_str
    if model == "pro":
        history = web_v4pro_history.get(user_id_int, web_v4pro_history.get(user_id_str, {}))
    else:
        history = web_v4flash_history.get(user_id_int, web_v4flash_history.get(user_id_str, {}))
    chat_history = history.get(chat_id, [])
    full_text = ""
    for msg in reversed(chat_history):
        if msg.get("role") == "assistant":
            full_text = msg["content"]
            break
    return jsonify({"full_text": full_text})

free_chat_history = {}
load_free_history()

SYSTEM_PROMPT = """Ты — администратор PhotoBot, без пола. Общаешься как живой человек. Твоя задача — помочь пользователю освоиться в сервисе, подобрать режим, составить промт для генерации, а также рассказать о DeepSeek. Если пользователь спрашивает кто ты — отвечаешь: «Я ИИ-ассистент PhotoBot на базе DeepSeek v4». Всегда общаешься вежливо и максимально расположительно к себе.

В PhotoBot есть два направления:
1. Генерация изображений на базе FLUX 2 — загрузка фото и создание новых образов
2. DeepSeek — текстовый ИИ для работы с кодом, текстами, документами, анализом данных

Как устроен PhotoBot (генерации фото):
Пользователь загружает фото через форму на сайте, выбирает режим, нейросеть обрабатывает около двух минут. Также есть режим Текст в изображение. Первые 5 генераций бесплатно, дальше — только по платной подписке. Результат сохраняется в Хранилище генераций. Эталонное фото загружается кнопкой "Использовать эталон", даёт самую большую вероятность сходства лица. Без эталона сходство может быть немного ниже ожидаемого. Одна подписка на все платформы.

Как устроен DeepSeek (текстовый ИИ):
Пользователь пополняет депозит и пользуется моделью DeepSeek v4. Доступны две модели:
• DeepSeek V4 Flash — быстрая, для кода, текстов, повседневных задач.
• DeepSeek V4 Pro — мощная, для сложных проектов, анализа, больших задач.
Пользователь платит только за использованные токены. Повторяющиеся части диалога кэшируются и стоят в 120 раз дешевле.

DeepSeek доступен на отдельных страницах сайта (V4 Flash и V4 Pro), а также в боте Telegram. Подписка единая. Пополнение через раздел подписки.

DeepSeek принимает вложения: .txt, .json, .csv, .xml, .html, .md, .py, .js, .java, .cpp, .c, .pdf, .docx, .xlsx, .pptx, .zip, .png, .jpg, .jpeg, .gif, .webp. Архивы .zip читает содержимое. Изображения распознаёт текст (OCR). PDF и офисные документы извлекает текст полностью.

Режимы генерации фото:
• Готовые (описание уже вшито, можно указать одежду): Документы, Закат на пляже, Цветы в авто, День рождения, Мужское в зеркале — рекомендуются для начала.
• Свободные: Свободный женский, Свободный мужской, Text2Image — когда пользователь освоился и готов к своему описанию.

Ты знаешь модель FLUX 2 klein 9b fp8. Пользователю говоришь только «FLUX 2». При составлении описаний используешь: камеры (Canon G7X Mark III, Sony A7IV, Fujifilm X-T5), линзы (85mm, 24mm, 50mm), освещение (золотой час, неон, студия), эффекты (плёнка Kodak Portra 400, Cinestill 800T), стили (фотореализм, кинематографичный, гиперреализм), композицию (правило третей, боке). Добавляй 2-3 технические детали в описание.

Твоя роль:
• Ты консультируешь пользователя по всем вопросам сервиса: как выбрать режим, как составить промт, как пользоваться DeepSeek, как оформить подписку.
• Помогаешь подобрать описание для генерации — уточняй детали (стиль, одежду, фон, свет, настроение), пока пользователь не подтвердит что всё сказал.
• Не генерируешь фото — только помогаешь советом. Генерацию пользователь запускает сам через форму на сайте.
• Если пользователь спрашивает про DeepSeek — расскажи о моделях V4 Flash и V4 Pro, форматах вложений, как пополнить баланс.

Правила общения:
• Общайся на «Вы», уважительно, тепло.
• Не пиши «понял», «принято», «выполняю».
• Используй один эмодзи по контексту: 🌅🌸👔🎂💪.
• База примеров: фото-бот.рф/prompts.
• Если общение ведется не по тематике PhotoBot или DeepSeek: «Помогаю только по вопросам сервиса PhotoBot».
• Обработка занимает около двух минут, не говори «быстро».
• Помни контекст: не здоровайся повторно, если ранее уже приветствовал пользователя.
• Твой ответ ограничен 400 токенами. Выстраивай формулировку ответа таким образом, чтобы корректно закончить мысль в пределах лимита. Если мысль не влезает в лимит — сократи её, скажи самое главное.
• Не рекомендуй пользователям чрезмерно длинные промты — они могут вызвать галлюцинации нейросети.
• Если пользователь 3 раза подряд отправляет одно и то же короткое сообщение (например «привет»), отвечай строго одной фразой: «Помогаю только по вопросам сервиса PhotoBot» и прекращай диалог."""

@app.route("/get_balance", methods=["GET"])
def get_balance():
    global user_v4pro_balance, user_v4flash_balance
    if not check_web_session(request):
        return jsonify({"error": "Unauthorized"}), 403
    user_id_str = request.args.get("user_id", "")
    user_id_int = int(user_id_str) if str(user_id_str).isdigit() else user_id_str
    if os.path.exists(V4PRO_BALANCE_FILE):
        with open(V4PRO_BALANCE_FILE, "rb") as f:
            user_v4pro_balance.update(pickle.load(f))
    if os.path.exists(V4FLASH_BALANCE_FILE):
        with open(V4FLASH_BALANCE_FILE, "rb") as f:
            user_v4flash_balance.update(pickle.load(f))
    pro = user_v4pro_balance.get(user_id_int, 0)
    flash = user_v4flash_balance.get(user_id_int, 0)
    return jsonify({"pro": pro, "flash": flash})

@app.route("/chat_free", methods=["POST"])
def chat_free():
    ip = request.headers.get("X-Forwarded-For", request.remote_addr)
    if not check_rate_limit(ip, "chat_free", 10, 60):
        return jsonify({"reply": "Слишком много запросов. Подождите минуту."})
    data = request.json
    user_id = data.get("user_id", "web_user")
    message = data.get("message", "")

    if user_id not in free_chat_history:
        free_chat_history[user_id] = [{"role": "system", "content": SYSTEM_PROMPT}]

    free_chat_history[user_id].append({"role": "user", "content": message})

    if len(free_chat_history[user_id]) > 51:
        free_chat_history[user_id] = free_chat_history[user_id][-51:]

    deepseek_api_key = os.getenv("DEEPSEEK_API_KEY")
    if not deepseek_api_key:
        return jsonify({"error": "API key not configured"}), 500

    headers = {
        "Authorization": f"Bearer {deepseek_api_key}",
        "Content-Type": "application/json"
    }

    payload = {
        "model": "deepseek-v4-flash",
        "messages": free_chat_history[user_id],
        "max_tokens": 400
    }

    try:
        resp = req.post("https://api.deepseek.com/v1/chat/completions", headers=headers, json=payload, timeout=60)
        if resp.status_code == 200:
            data = resp.json()
            reply = data["choices"][0]["message"]["content"]
            free_chat_history[user_id].append({"role": "assistant", "content": reply})
            safe_pickle_dump(free_chat_history, FREE_CHAT_HISTORY_FILE)
            return jsonify({"reply": reply})
        else:
            return jsonify({"error": "API error"}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/get_free_chat_history", methods=["GET"])
def get_free_chat_history():
    user_id = request.args.get("user_id", "")
    history = free_chat_history.get(user_id, [])
    return jsonify({"history": history})

@app.route("/get_chat_history", methods=["GET"])
def get_chat_history():
    if not check_web_session(request):
        return jsonify({"error": "Unauthorized"}), 403
    user_id_str = request.args.get("user_id", "")
    model = request.args.get("model", "flash")
    user_id_int = int(user_id_str) if user_id_str.isdigit() else user_id_str
    if model == "pro":
        history = web_v4pro_history.get(user_id_int, web_v4pro_history.get(user_id_str, {}))
    else:
        history = web_v4flash_history.get(user_id_int, web_v4flash_history.get(user_id_str, {}))
    return jsonify({"history": history})

@app.route("/chat", methods=["POST"])
def chat():
    if not check_web_session(request):
        return jsonify({"error": "Unauthorized"}), 403

    data = request.json
    user_id = data.get("user_id", "web_user")
    message = data.get("message", "")
    model = data.get("model", "flash")
    mode = data.get("mode", "chat")
    file_data = data.get("file_data", "")
    file_name = data.get("file_name", "")
    file_mime = data.get("file_mime", "")
    chat_id = data.get("chat_id", "default")

    deepseek_api_key = os.getenv("DEEPSEEK_API_KEY")
    if not deepseek_api_key:
        return jsonify({"error": "API key not configured"}), 500

    thinking = {"type": "enabled"} if mode == "reasoner" else {"type": "disabled"}

    if model == "pro":
        model_name = "deepseek-v4-pro"
    else:
        model_name = "deepseek-v4-flash"

    headers = {
        "Authorization": f"Bearer {deepseek_api_key}",
        "Content-Type": "application/json"
    }
    global user_v4pro_balance, user_v4flash_balance

    user_msg = {"role": "user", "content": message or ""}
    if file_data:
        print(f"[CHAT] file_name={file_name} file_mime={file_mime} len={len(file_data)}")
        if file_mime and file_mime.startswith("image/"):
            try:
                img_bytes = base64.b64decode(file_data)
                img = Image.open(io.BytesIO(img_bytes))
                ocr_text = pytesseract.image_to_string(img, lang='rus+eng').strip()
                if ocr_text:
                    if message:
                        user_msg["content"] = f"{message}\n\n--- Текст с изображения ({file_name}) ---\n{ocr_text}"
                    else:
                        user_msg["content"] = f"Текст с изображения ({file_name}):\n{ocr_text}"
                else:
                    if message:
                        user_msg["content"] = f"{message}\n\n[Изображение без текста: {file_name}]"
                    else:
                        user_msg["content"] = f"[Изображение без текста: {file_name}]"
            except Exception as e:
                print(f"[OCR] Error: {e}")
                if message:
                    user_msg["content"] = f"{message}\n\n[Ошибка распознавания изображения: {file_name}]"
                else:
                    user_msg["content"] = f"[Ошибка распознавания изображения: {file_name}]"
        else:
            ext = os.path.splitext(file_name)[1].lower() if file_name else ''
            file_bytes = base64.b64decode(file_data)
            text_content = ''
            try:
                if ext == '.pdf':
                    try:
                        import fitz
                        doc = fitz.open(stream=file_bytes, filetype='pdf')
                        pages = []
                        for page in doc:
                            pages.append(page.get_text())
                        doc.close()
                        text_content = '\n'.join(pages)
                    except ImportError:
                        text_content = "[PDF парсер не установлен. Установите PyMuPDF]"
                elif ext == '.docx':
                    try:
                        from docx import Document
                        doc = Document(io.BytesIO(file_bytes))
                        text_content = '\n'.join([p.text for p in doc.paragraphs])
                    except ImportError:
                        text_content = "[DOCX парсер не установлен. Установите python-docx]"
                elif ext == '.xlsx':
                    try:
                        from openpyxl import load_workbook
                        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
                        rows = []
                        for sheet_name in wb.sheetnames:
                            ws = wb[sheet_name]
                            rows.append(f'--- {sheet_name} ---')
                            for row in ws.iter_rows(values_only=True):
                                rows.append('\t'.join([str(c) if c is not None else '' for c in row]))
                        wb.close()
                        text_content = '\n'.join(rows)
                    except ImportError:
                        text_content = "[XLSX парсер не установлен. Установите openpyxl]"
                elif ext == '.pptx':
                    try:
                        from pptx import Presentation
                        prs = Presentation(io.BytesIO(file_bytes))
                        slides = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, 'text') and shape.text.strip():
                                    slides.append(shape.text.strip())
                        text_content = '\n'.join(slides)
                    except ImportError:
                        text_content = "[PPTX парсер не установлен. Установите python-pptx]"
                elif ext == '.zip':
                    import zipfile
                    zf = zipfile.ZipFile(io.BytesIO(file_bytes))
                    items = []
                    for name in zf.namelist():
                        if not name.endswith('/'):
                            try:
                                content = zf.read(name).decode('utf-8', errors='replace')
                                items.append(f'--- {name} ---\n{content}')
                            except:
                                items.append(f'--- {name} ---\n[Бинарный файл, не читается]')
                    zf.close()
                    text_content = '\n'.join(items)
                    if len(text_content) > 100000:
                        text_content = text_content[:100000] + '\n\n[Содержимое архива обрезано — слишком большой объём]'
                else:
                    text_content = file_bytes.decode('utf-8', errors='replace')
            except Exception as e:
                print(f"[FILE] Parse error for {ext}: {e}")
                try:
                    text_content = file_bytes.decode('utf-8', errors='replace')
                except:
                    text_content = "[Файл не может быть отображён как текст]"
            if not text_content.strip():
                text_content = "[Файл не содержит текста]"
            if message:
                user_msg["content"] = f"{message}\n\n--- Содержимое файла ({file_name}) ---\n{text_content}"
            else:
                user_msg["content"] = f"Содержимое файла ({file_name}):\n{text_content}"

    if model == "pro":
        if user_id not in web_v4pro_history:
            web_v4pro_history[user_id] = {}
        if chat_id not in web_v4pro_history[user_id]:
            web_v4pro_history[user_id][chat_id] = []
        history = web_v4pro_history[user_id][chat_id]
    else:
        if user_id not in web_v4flash_history:
            web_v4flash_history[user_id] = {}
        if chat_id not in web_v4flash_history[user_id]:
            web_v4flash_history[user_id][chat_id] = []
        history = web_v4flash_history[user_id][chat_id]

    # Если история пуста — добавляем системный промпт
    if len(history) == 0:
        system_prompt_text = "Ты DeepSeek V4 Flash. Отвечай обычным текстом. Не ставь #, *, ** в тексте. Не используй заголовки, списки с тире, жирный шрифт, курсив. Пиши простым текстом с переносами строк."
        if model == "pro":
            system_prompt_text = system_prompt_text.replace("V4 Flash", "V4 Pro")
        history.append({"role": "system", "content": system_prompt_text})

    history.append(user_msg)
    messages = list(history)

    # Оценка стоимости запроса
    if model == "pro":
        max_cost = 384000 * 0.00000348 * 48 * 1.3
        current_balance = user_v4pro_balance.get(user_id, 0)
    else:
        max_cost = 384000 * 0.00000028 * 48 * 1.3
        current_balance = user_v4flash_balance.get(user_id, 0)

    if current_balance <= 0:
        return jsonify({"error": "insufficient_balance"}), 402

    if current_balance < max_cost:
        return jsonify({"error": "insufficient_balance", "required": round(max_cost, 2), "balance": round(current_balance, 2)}), 402

    payload = {
        "model": model_name,
        "messages": messages,
        "max_tokens": 384000
    }
    if thinking:
        payload["thinking"] = thinking

    try:
        resp = req.post("https://api.deepseek.com/v1/chat/completions", headers=headers, json=payload, timeout=60)
        if resp.status_code == 200:
            data = resp.json()
            msg = data["choices"][0]["message"]
            reply = msg["content"]
            reasoning = msg.get("reasoning_content", "")

            if os.path.exists(V4PRO_BALANCE_FILE):
                with open(V4PRO_BALANCE_FILE, "rb") as f:
                    user_v4pro_balance.update(pickle.load(f))
            if os.path.exists(V4FLASH_BALANCE_FILE):
                with open(V4FLASH_BALANCE_FILE, "rb") as f:
                    user_v4flash_balance.update(pickle.load(f))
            usage = data.get("usage", {})
            prompt_tokens = usage.get("prompt_tokens", 0)
            completion_tokens = usage.get("completion_tokens", 0)
            reasoning_tokens = usage.get("reasoning_tokens", 0) or usage.get("completion_tokens_details", {}).get("reasoning_tokens", 0)
            if model == "pro":
                web_v4pro_history[user_id][chat_id].append({"role": "assistant", "content": reply})
                safe_pickle_dump(web_v4pro_history, WEB_V4PRO_HISTORY_FILE)
            else:
                web_v4flash_history[user_id][chat_id].append({"role": "assistant", "content": reply})
                safe_pickle_dump(web_v4flash_history, WEB_V4FLASH_HISTORY_FILE)
            if model == "pro":
                cache_hit = usage.get("prompt_cache_hit_tokens", 0)
                cost = ((prompt_tokens - cache_hit) * 0.00000174 + cache_hit * 0.00000174 / 120 + (completion_tokens + reasoning_tokens) * 0.00000348) * 48 * 1.3
                if user_v4pro_balance.get(user_id, 0) >= cost:
                    user_v4pro_balance[user_id] = user_v4pro_balance.get(user_id, 0) - cost
                else:
                    user_v4pro_balance[user_id] = 0
            else:
                cache_hit = usage.get("prompt_cache_hit_tokens", 0)
                cost = ((prompt_tokens - cache_hit) * 0.00000014 + cache_hit * 0.00000014 / 120 + (completion_tokens + reasoning_tokens) * 0.00000028) * 48 * 1.3
                if user_v4flash_balance.get(user_id, 0) >= cost:
                    user_v4flash_balance[user_id] = user_v4flash_balance.get(user_id, 0) - cost
                else:
                    user_v4flash_balance[user_id] = 0

            safe_pickle_dump(user_v4pro_balance, V4PRO_BALANCE_FILE)
            safe_pickle_dump(user_v4flash_balance, V4FLASH_BALANCE_FILE)
            return jsonify({"reply": reply, "reasoning": reasoning})
        else:
            print(f"[CHAT] Error: {resp.status_code} {resp.text}")
            return jsonify({"error": "API error"}), 500
    except Exception as e:
        print(f"[CHAT] Exception: {e}")
        return jsonify({"error": str(e)}), 500

def process_app_queue():
    print("[QUEUE] Worker started")
    while True:
        time.sleep(2)
        with app_jobs_lock:
            pending_jobs = [(jid, j) for jid, j in app_jobs.items() if j.get("status") == "pending"]
        print(f"[QUEUE] Pending jobs: {len(pending_jobs)}")
        for job_id, job in pending_jobs:
            try:
                mode = job.get("mode", "documents")
                caption = job.get("caption", "")
                ref_name = job.get("ref_name", "")
                img_data = job.get("image_data")
                image_name = None
                if mode != "wallpaper":
                    if ref_name:
                        image_name = ref_name
                    elif img_data:
                        image_name = f"app_{uuid.uuid4().hex[:8]}.jpg"
                        save_image_to_input(img_data, image_name)
                    else:
                        with app_jobs_lock:
                            app_jobs[job_id]["status"] = "failed"
                        continue
                caption_val = caption.strip().replace('\n', ' ').replace('\r', ' ').replace('"', "'") if caption.strip() else ""
                if mode == "documents":
                    caption_val = caption_val if caption_val else "черная рубашка"
                elif mode == "birthday":
                    parts = caption_val.split() if caption_val else []
                    caption_val = " ".join(parts) if parts else "23 красное"
                elif mode == "flowers":
                    caption_val = caption_val if caption_val else "белое"
                wf = load_workflow(mode, image_name, caption_val)
                if wf is None:
                    print(f"[QUEUE] Workflow not found for {job_id}")
                    with app_jobs_lock:
                        app_jobs[job_id]["status"] = "failed"
                    continue
                prompt_resp = req.post(f"{COMFYUI_URL}/prompt", json={"prompt": wf}).json()
                prompt_id = prompt_resp.get("prompt_id")
                if not prompt_id:
                    print(f"[QUEUE] Prompt failed for {job_id}")
                    with app_jobs_lock:
                        app_jobs[job_id]["status"] = "failed"
                    continue
                print(f"[QUEUE] Prompt {prompt_id} for {job_id}")
                found = False
                for i in range(180):
                    time.sleep(1)
                    history = req.get(f"{COMFYUI_URL}/history").json()
                    if prompt_id in history:
                        outputs = history[prompt_id]["outputs"]
                        for nid in outputs:
                            for img in outputs[nid].get("images", []):
                                if img["filename"].startswith("Flux2-Klein"):
                                    img_resp = req.get(f"{COMFYUI_URL}/view?filename={img['filename']}&subfolder={img.get('subfolder','')}&type=output")
                                    out_path = os.path.join(TEMP_FOLDER, f"app_{job_id}.jpg")
                                    with open(out_path, "wb") as f:
                                        f.write(img_resp.content)
                                    web_path = os.path.join(DATASET_FOLDER, f"web_{job['user_id']}_{job_id}.jpg")
                                    with open(out_path, "rb") as src, open(web_path, "wb") as dst:
                                        dst.write(src.read())
                                    txt_path = os.path.join(DATASET_FOLDER, f"web_{job['user_id']}_{job_id}.txt")
                                    with open(txt_path, "w", encoding="utf-8") as tf:
                                        tf.write(caption_val)
                                    with app_jobs_lock:
                                        app_jobs[job_id]["status"] = "ready"
                                        app_jobs[job_id]["file"] = out_path
                                    print(f"[QUEUE] Completed {job_id}")
                                    send_push(job["user_id"], "Генерация завершена", "Ваше фото готово!")
                                    found = True
                                    if job.get("source") == "web":
                                        try:
                                            user_id_int = int(job["user_id"]) if str(job["user_id"]).isdigit() else job["user_id"]
                                            free_generations[user_id_int] = free_generations.get(user_id_int, 0) + 1
                                            safe_pickle_dump(free_generations, FREE_GENERATIONS_FILE)
                                            print(f"[QUEUE] Free generation #{free_generations[user_id_int]} for user {user_id_int}")
                                        except Exception as e:
                                            print(f"[QUEUE] Error updating free_generations: {e}")

                                    break
                            if found:
                                break
                        if found:
                            break
                if not found:
                    with app_jobs_lock:
                        app_jobs[job_id]["status"] = "failed"
            except Exception as e:
                print(f"[QUEUE] Error for {job_id}: {e}")
                with app_jobs_lock:
                    app_jobs[job_id]["status"] = "failed"

@app.route("/delete_history_file", methods=["DELETE"])
def delete_history_file():
    source = request.args.get("source", "app")
    if source == 'web':
        if not check_web_session(request):
            return jsonify({"error": "Unauthorized"}), 403
    elif not check_api_token(request):
        return jsonify({"error": "Unauthorized"}), 403
    filename = request.args.get("file", "")
    path = os.path.join(DATASET_FOLDER, filename)
    if os.path.exists(path):
        os.remove(path)
        txt_path = path.replace('.jpg', '.txt')
        if os.path.exists(txt_path):
            os.remove(txt_path)
        return jsonify({"success": True})
    return jsonify({"error": "Not found"}), 404

@app.route("/create_payment", methods=["POST"])
def create_payment():
    if not YOOKASSA_SHOP_ID or not YOOKASSA_SECRET_KEY:
        return jsonify({"error": "Payment system not configured"}), 500
    ip = request.headers.get("X-Forwarded-For", request.remote_addr)
    if not check_rate_limit(ip, "create_payment", 5, 3600):
        return jsonify({"error": "Слишком много попыток. Попробуйте позже."})
    if not check_api_token(request):
        return jsonify({"error": "Unauthorized"}), 403
    data = request.json
    user_id = data.get("user_id")
    amount_rub = data.get("amount_rub")
    description = data.get("description", "Пополнение PhotoBot")

    idempotence_key = uuid.uuid4().hex

    payload = {
        "amount": {
            "value": str(amount_rub),
            "currency": "RUB"
        },
        "confirmation": {
            "type": "redirect",
            "return_url": "https://t.me/PhotoDocumentsBot_bot"
        },
        "description": description,
        "metadata": {
            "user_id": str(user_id)
        }
    }

    try:
        resp = req.post(
            "https://api.yookassa.ru/v3/payments",
            auth=(YOOKASSA_SHOP_ID, YOOKASSA_SECRET_KEY),
            json=payload,
            headers={"Idempotence-Key": idempotence_key},
            timeout=30
        )
        if resp.status_code == 200:
            payment_data = resp.json()
            return jsonify({
                "success": True,
                "payment_url": payment_data["confirmation"]["confirmation_url"],
                "payment_id": payment_data["id"]
            })
        else:
            print(f"[YOOKASSA] Error: {resp.status_code} {resp.text}")
            return jsonify({"error": "Payment creation failed"}), 500
    except Exception as e:
        print(f"[YOOKASSA] Exception: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/yookassa_webhook", methods=["POST"])
def yookassa_webhook():
    data = request.json
    if data.get("event") == "payment.succeeded":
        payment = data.get("object", {})
        amount = payment.get("amount", {}).get("value", "0")
        user_id = payment.get("metadata", {}).get("user_id", "")
        description = payment.get("description", "")

        print(f"[YOOKASSA] Payment succeeded: {amount} RUB from user {user_id}")

        amount_rub = float(amount)
        user_id_int = int(user_id) if user_id.isdigit() else user_id

        if "PhotoBot" in description:
            VIP_USERS.add(user_id_int)
            if "month" in description:
                VIP_EXPIRY[user_id_int] = datetime.now() + timedelta(days=30)
            elif "year" in description:
                VIP_EXPIRY[user_id_int] = datetime.now() + timedelta(days=365)
            safe_pickle_dump(VIP_EXPIRY, VIP_EXPIRY_FILE)
            save_vip()
        elif "V4 Pro" in description:
            user_v4pro_balance[user_id_int] = user_v4pro_balance.get(user_id_int, 0) + amount_rub * 0.8
            safe_pickle_dump(user_v4pro_balance, V4PRO_BALANCE_FILE)
        elif "V4 Flash" in description:
            user_v4flash_balance[user_id_int] = user_v4flash_balance.get(user_id_int, 0) + amount_rub * 1.5
            safe_pickle_dump(user_v4flash_balance, V4FLASH_BALANCE_FILE)

        return jsonify({"success": True})
    return jsonify({"success": True})

@app.route("/check_generation_access", methods=["GET"])
def check_generation_access():
    user_id_str = request.args.get("user_id", "")
    user_id_int = int(user_id_str) if user_id_str.isdigit() else user_id_str
    # Загружаем свежие данные
    if os.path.exists(FREE_GENERATIONS_FILE):
        with open(FREE_GENERATIONS_FILE, "rb") as f:
            free_generations.update(pickle.load(f))
    if os.path.exists(VIP_FILE):
        with open(VIP_FILE, "rb") as f:
            VIP_USERS.update(set(pickle.load(f)))
    gen_count = free_generations.get(user_id_int, 0)
    if gen_count < 5:
        return jsonify({"allowed": True, "free": True, "remaining": 5 - gen_count})
    if user_id_int in VIP_USERS:
        if user_id_int in VIP_EXPIRY and datetime.now() > VIP_EXPIRY[user_id_int]:
            pass  # истекла — идём дальше к отказу
        else:
            return jsonify({"allowed": True, "free": False})
    return jsonify({"allowed": False, "free": False, "remaining": 0})

@app.route("/delete_chat", methods=["DELETE"])
def delete_chat():
    if not check_web_session(request):
        return jsonify({"error": "Unauthorized"}), 403
    user_id_str = request.args.get("user_id", "")
    model = request.args.get("model", "flash")
    chat_id = request.args.get("chat_id", "")
    user_id_int = int(user_id_str) if user_id_str.isdigit() else user_id_str

    if model == "pro":
        if user_id_int in web_v4pro_history and chat_id in web_v4pro_history[user_id_int]:
            del web_v4pro_history[user_id_int][chat_id]
            safe_pickle_dump(web_v4pro_history, WEB_V4PRO_HISTORY_FILE)
    else:
        if user_id_int in web_v4flash_history and chat_id in web_v4flash_history[user_id_int]:
            del web_v4flash_history[user_id_int][chat_id]
            safe_pickle_dump(web_v4flash_history, WEB_V4FLASH_HISTORY_FILE)
    return jsonify({"success": True})

if __name__ == "__main__":
    print("[QUEUE] Starting worker thread...")
    threading.Thread(target=process_app_queue, daemon=True).start()
    app.run(host="0.0.0.0", port=5000, debug=False)
